#!/usr/bin/env python3
"""生成「和 AI 作朋友——高中生的 AI 工具箱」工作坊簡報 .pptx

對齊 slides.html 的 38 張投影片結構，並從 forTeacher.org 抽取教學情境，
為每張投影片加入演講者備忘稿（speaker notes）。

執行：python3 generate_slides.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── 配色（對齊 slides-styles.css 主題） ──
BG_DARK = RGBColor(0x1A, 0x1A, 0x2E)
BG_SECTION = RGBColor(0x16, 0x21, 0x3E)
CARD_BG = RGBColor(0x22, 0x2B, 0x45)
ACCENT = RGBColor(0x00, 0xD2, 0xFF)         # cyan
ACCENT2 = RGBColor(0xFF, 0x6B, 0x6B)        # coral / red
ACCENT3 = RGBColor(0x4E, 0xCB, 0x71)        # green
ACCENT4 = RGBColor(0xFF, 0xAA, 0x33)        # orange
ACCENT5 = RGBColor(0xCC, 0x66, 0xFF)        # purple
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xBB, 0xBB, 0xCC)
SUBTITLE_GRAY = RGBColor(0x99, 0x99, 0xAA)

FONT_TITLE = "Helvetica Neue"
FONT_BODY = "Helvetica Neue"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SLIDE_W = prs.slide_width
SLIDE_H = prs.slide_height


# ═══════════════════════════════════════════
# 共用 helpers
# ═══════════════════════════════════════════

def new_slide(bg=BG_DARK):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = bg
    return slide


def add_rect(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_text(slide, left, top, width, height, text,
             size=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT,
             font=FONT_BODY, line_spacing=1.4):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Pt(0)
    tf.margin_top = tf.margin_bottom = Pt(0)
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font
    p.alignment = align
    p.line_spacing = line_spacing
    return box


def add_bullets(slide, left, top, width, height, items,
                size=22, color=LIGHT_GRAY, bullet_color=ACCENT, spacing=1.6):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Pt(0)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        bullet = p.add_run()
        bullet.text = "▸  "
        bullet.font.size = Pt(size)
        bullet.font.color.rgb = bullet_color
        bullet.font.name = FONT_BODY
        bullet.font.bold = True
        body = p.add_run()
        body.text = item
        body.font.size = Pt(size)
        body.font.color.rgb = color
        body.font.name = FONT_BODY
        p.line_spacing = spacing
        p.space_after = Pt(size * 0.4)
    return box


def add_callout(slide, left, top, width, height, text,
                color=ACCENT, label=None, size=16):
    """色塊 callout：左側色條 + 半透明底色"""
    add_rect(slide, left, top, width, height, CARD_BG)
    add_rect(slide, left, top, Pt(5), height, color)
    if label:
        add_text(slide, left + Inches(0.2), top + Inches(0.1), width - Inches(0.4), Inches(0.4),
                 label, size=14, color=color, bold=True)
        add_text(slide, left + Inches(0.2), top + Inches(0.55), width - Inches(0.4),
                 height - Inches(0.6), text, size=size, color=WHITE)
    else:
        add_text(slide, left + Inches(0.2), top + Inches(0.2), width - Inches(0.4),
                 height - Inches(0.4), text, size=size, color=WHITE)


def add_notes(slide, text):
    """加入演講者備忘稿"""
    notes_tf = slide.notes_slide.notes_text_frame
    notes_tf.text = text
    return notes_tf


# ═══════════════════════════════════════════
# Slide patterns
# ═══════════════════════════════════════════

def title_slide():
    slide = new_slide(BG_SECTION)
    add_rect(slide, Inches(5.67), Inches(1.6), Inches(2), Pt(3), ACCENT)
    add_text(slide, Inches(0), Inches(1.9), SLIDE_W, Inches(1.4),
             "和 AI 作朋友", size=64, bold=True, color=WHITE,
             font=FONT_TITLE, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(0), Inches(3.4), SLIDE_W, Inches(0.8),
             "高中生的 AI 工具箱", size=32, color=ACCENT,
             font=FONT_TITLE, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(0), Inches(4.5), SLIDE_W, Inches(0.5),
             "顏永進  ｜  2026.04.30  ｜  雲林", size=20,
             color=SUBTITLE_GRAY, align=PP_ALIGN.CENTER)
    for x, label, url in [
        (Inches(3.5), "簡報", "is.gd/1rLH3E"),
        (Inches(7.8), "教學網頁", "is.gd/tENFQN"),
    ]:
        add_rect(slide, x, Inches(5.4), Inches(1.5), Inches(1.5), CARD_BG)
        add_text(slide, x, Inches(6.0), Inches(1.5), Inches(0.4),
                 "QR", size=18, color=SUBTITLE_GRAY, align=PP_ALIGN.CENTER)
        add_text(slide, x, Inches(7.0), Inches(1.5), Inches(0.3),
                 label, size=14, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
        add_text(slide, x + Inches(1.6), Inches(5.7), Inches(2.5), Inches(0.4),
                 url, size=18, color=ACCENT, bold=True)
    return slide


def section_divider(label_top, title, subtitle, color=ACCENT):
    slide = new_slide(BG_SECTION)
    add_rect(slide, Inches(0), Inches(0), Inches(0.4), SLIDE_H, color)
    add_text(slide, Inches(1.2), Inches(2.5), Inches(11), Inches(0.6),
             label_top, size=20, color=color, bold=True)
    add_text(slide, Inches(1.2), Inches(3.2), Inches(11), Inches(1.4),
             title, size=54, bold=True, color=WHITE, font=FONT_TITLE)
    add_text(slide, Inches(1.2), Inches(4.7), Inches(11), Inches(0.6),
             subtitle, size=22, color=LIGHT_GRAY)
    return slide


def step_divider(step_num, tool_name, subtitle, color=ACCENT):
    slide = new_slide(BG_SECTION)
    add_rect(slide, Inches(0), Inches(0), Inches(0.4), SLIDE_H, color)
    add_text(slide, Inches(1.2), Inches(2.4), Inches(11), Inches(0.5),
             f"STEP {step_num}", size=20, color=color, bold=True)
    add_text(slide, Inches(1.2), Inches(3.0), Inches(11), Inches(1.6),
             tool_name, size=66, bold=True, color=WHITE, font=FONT_TITLE)
    add_text(slide, Inches(1.2), Inches(4.8), Inches(11), Inches(0.6),
             subtitle, size=24, color=LIGHT_GRAY)
    return slide


def content_slide(step_label, title, color=ACCENT):
    slide = new_slide(BG_DARK)
    if step_label:
        badge = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(0.8), Inches(0.6), Inches(0.55), Inches(0.55))
        badge.fill.solid()
        badge.fill.fore_color.rgb = color
        badge.line.fill.background()
        tf = badge.text_frame
        tf.margin_left = tf.margin_right = Pt(0)
        tf.margin_top = tf.margin_bottom = Pt(0)
        p = tf.paragraphs[0]
        p.text = str(step_label)
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    add_text(slide, Inches(1.6) if step_label else Inches(0.8), Inches(0.65),
             Inches(11), Inches(0.7),
             title, size=32, bold=True, color=WHITE, font=FONT_TITLE)
    add_rect(slide, Inches(0.8), Inches(1.45), Inches(11.5), Pt(2), color)
    return slide


# ═══════════════════════════════════════════
# 1. 封面
# ═══════════════════════════════════════════
slide = title_slide()
add_notes(slide, """開場 30 秒。

「AI 不是未來的事，是現在的事。問題不是要不要用，而是你會不會用。」
還在 prompt → copy → paste？今天帶你看真正的玩法。

提醒學生兩個 QR code：上面是這份簡報，下面是教學網頁（含完整文字、素材清單、附錄）。""")


# ═══════════════════════════════════════════
# 2. 今天要做什麼？
# ═══════════════════════════════════════════
slide = content_slide(None, "今天要做什麼？", ACCENT)
add_text(slide, Inches(0.8), Inches(1.75), Inches(11.5), Inches(0.6),
         "很多人以為 AI 是「幫你寫東西」——真正厲害的是「讓 AI 幫你做事情」",
         size=18, color=LIGHT_GRAY)
add_text(slide, Inches(0.8), Inches(2.35), Inches(11.5), Inches(0.5),
         "三小時練「交辦」，做出一份真的可以上傳的台灣史學習歷程",
         size=18, color=LIGHT_GRAY)
add_text(slide, Inches(0.8), Inches(2.95), Inches(11.5), Inches(0.6),
         "主題：為什麼日本人選雲林蓋糖廠？",
         size=24, bold=True, color=ACCENT, font=FONT_TITLE)

stages = [
    ("讀", "NotebookLM 消化糖業史料", "含出處的研究筆記（Word）", ACCENT),
    ("做", "Napkin 畫概念圖 + Claude 生 Word", "小論文 .docx 初稿", ACCENT2),
    ("學", "ChatGPT 當教練改反思", "有深度的反思，回 Word", ACCENT3),
    ("簡報", "Gamma 或 Claude 做投影片", "學習歷程簡報", ACCENT4),
    ("創造", "（加分）AI Agent 做網頁", "糖業史時間軸網站", ACCENT5),
]
for i, (stage, what, output, color) in enumerate(stages):
    y = Inches(3.85 + i * 0.65)
    add_rect(slide, Inches(0.8), y, Inches(11.5), Inches(0.55), CARD_BG)
    add_rect(slide, Inches(0.8), y, Pt(5), Inches(0.55), color)
    add_text(slide, Inches(1.0), y + Inches(0.1), Inches(0.9), Inches(0.4),
             stage, size=18, bold=True, color=color, font=FONT_TITLE)
    add_text(slide, Inches(2.0), y + Inches(0.13), Inches(5.8), Inches(0.4),
             what, size=15, color=WHITE)
    add_text(slide, Inches(7.9), y + Inches(0.13), Inches(4.4), Inches(0.4),
             output, size=14, color=LIGHT_GRAY)

add_notes(slide, """主線：用 AI 做一份雲林糖業的台灣史學習歷程。

關鍵訊息：
- 不是 AI 幫你「寫東西」，是你會不會用 AI「做事情」（交辦）
- 三小時跑完五階段：讀→做→學→簡報→創造
- 每階段產出餵進下一階段——最後手上有一份 Word + 一份簡報，可能還有互動網頁

為什麼選這個主題：學生人在雲林、虎尾糖廠就在隔壁；108 課綱第一冊第五章；在地題材撞題率低。""")


# ═══════════════════════════════════════════
# 3. AI 五層次
# ═══════════════════════════════════════════
slide = content_slide(None, "AI 使用的五個層次", ACCENT)
add_text(slide, Inches(0.8), Inches(1.85), Inches(11.5), Inches(0.5),
         "大部分人卡在 Lv.1，所以寫出來的東西又累又像罐頭",
         size=18, color=LIGHT_GRAY)

levels = [
    ("Lv.1", "問答機", "prompt → copy → paste", "ChatGPT 網頁版", ACCENT2),
    ("Lv.2", "整合工作流", "AI 直接嵌在你的工具裡", "Copilot、NotebookLM", ACCENT4),
    ("Lv.3", "AI Agent（初級助理）", "AI 自己規劃步驟、操作電腦", "Claude Code、Gemini CLI", ACCENT),
    ("Lv.4", "進階助理", "你規劃步驟，AI 幫你操作電腦", "Claude Code / CLIo", ACCENT3),
    ("Lv.5", "專屬助理", "AI 幫你規劃步驟，幫你操作電腦", "OpenClaw、Hermes", ACCENT5),
]
for i, (lv, name, how, tools, color) in enumerate(levels):
    y = Inches(2.6 + i * 0.78)
    is_today = lv == "Lv.3"
    bg = RGBColor(0x2E, 0x40, 0x60) if is_today else CARD_BG
    add_rect(slide, Inches(0.8), y, Inches(11.5), Inches(0.65), bg)
    add_rect(slide, Inches(0.8), y, Pt(5), Inches(0.65), color)
    add_text(slide, Inches(1.0), y + Inches(0.13), Inches(0.9), Inches(0.4),
             lv, size=18, bold=True, color=color, font=FONT_TITLE)
    add_text(slide, Inches(2.0), y + Inches(0.13), Inches(2.0), Inches(0.4),
             name, size=18, bold=True, color=WHITE, font=FONT_TITLE)
    add_text(slide, Inches(4.2), y + Inches(0.18), Inches(4.5), Inches(0.35),
             how, size=14, color=LIGHT_GRAY)
    add_text(slide, Inches(8.8), y + Inches(0.18), Inches(3.5), Inches(0.35),
             tools, size=14, color=SUBTITLE_GRAY, align=PP_ALIGN.RIGHT)
    if is_today:
        add_text(slide, Inches(11.95), y + Inches(0.13), Inches(0.4), Inches(0.4),
                 "★", size=22, bold=True, color=ACCENT, font=FONT_TITLE)

add_text(slide, Inches(0.8), Inches(6.65), Inches(11.5), Inches(0.4),
         "今天直接跳到 Lv.3，並在收尾「點到」Lv.4——你要學的不是工具，是怎麼跟 AI 一起長",
         size=14, color=ACCENT3)

add_notes(slide, """重點：學生大多卡在 Lv.1（問一句、複製一段），所以寫出來像罐頭。

今天直接跳到 Lv.3（AI Agent），收尾「點到」Lv.4（Skills/MCP）讓他們知道學習不是停在這場工作坊。

Lv.5 (OpenClaw, Hermes) 是 2026 在進入主流的「專屬助理」概念，提醒：你要學的是「怎麼跟 AI 一起長」，不是某個工具。""")


# ═══════════════════════════════════════════
# Step 1: NotebookLM
# ═══════════════════════════════════════════
slide = step_divider(1, "NotebookLM", "餵素材、問問題、產出研究筆記", ACCENT)
add_notes(slide, """第一段「讀」20 min。

場景：你想寫一份台灣史的課程學習成果，但三份史料根本還沒消化。

接下來四張投影片：先試 ChatGPT 看 Lv.1 限制 → NotebookLM 怎麼用 → 打包指令 → 匯出 Word。""")


# Step 1 - 先看 Lv.1 限制
slide = content_slide("1", "先看 Lv.1 的限制", ACCENT)
add_text(slide, Inches(0.8), Inches(1.85), Inches(11.5), Inches(0.5),
         "打開 ChatGPT，直接問：", size=20, color=LIGHT_GRAY)
add_callout(slide, Inches(0.8), Inches(2.5), Inches(11.5), Inches(1.0),
            "「日本人為什麼選擇在雲林虎尾蓋糖廠？」",
            color=ACCENT2, size=18)
add_text(slide, Inches(0.8), Inches(3.85), Inches(11.5), Inches(0.4),
         "ChatGPT 給你看起來不錯的答案，但：",
         size=18, color=LIGHT_GRAY)
add_bullets(slide, Inches(1.2), Inches(4.4), Inches(11.0), Inches(2.2), [
    "可能混入不正確的細節（幻覺）",
    "不會提到「戊戌大水災」這個關鍵事件",
    "沒有引用來源——你不知道哪句是真的",
], size=18, bullet_color=ACCENT2)
add_callout(slide, Inches(0.8), Inches(6.5), Inches(11.5), Inches(0.65),
            "這就是 Lv.1 的極限——給答案，但無法驗證。換 Lv.2。",
            color=ACCENT4, size=15)

add_notes(slide, """現場示範：直接打開 chatgpt.com，貼上這個問題。

ChatGPT 會回一個看起來合理的答案（氣候、土地之類），但：
- 細節可能是幻覺
- 通常不會提到「戊戌大水災」這個關鍵事件
- 完全沒有出處

對話結束時下一句：「這就是 Lv.1——你問一個問題，它給你一個答案，但你沒辦法驗證。接下來換 Lv.2 的方式。」""")


# Step 1 - NotebookLM 怎麼用
slide = content_slide("1", "NotebookLM 怎麼用？", ACCENT)
add_text(slide, Inches(0.8), Inches(1.85), Inches(5.5), Inches(0.5),
         "上傳素材", size=22, bold=True, color=ACCENT, font=FONT_TITLE)
add_text(slide, Inches(0.8), Inches(2.5), Inches(5.5), Inches(2.2),
         "把四份糖業史料 PDF 上傳到 NotebookLM\n\n素材清單見教學網頁底部",
         size=17, color=LIGHT_GRAY, line_spacing=1.5)
add_text(slide, Inches(7.0), Inches(1.85), Inches(5.7), Inches(0.5),
         "依序問四個問題", size=22, bold=True, color=ACCENT, font=FONT_TITLE)
add_bullets(slide, Inches(7.0), Inches(2.5), Inches(5.7), Inches(4), [
    "大方向：為什麼選雲林虎尾？",
    "追細節：戊戌大水災跟糖廠的關係？",
    "比較：虎尾 vs. 橋頭有什麼不同？",
    "切入：「為什麼是雲林？」大綱怎麼寫？",
], size=16, bullet_color=ACCENT)

add_notes(slide, """上傳 4 份糖業史料 PDF（事前已準備）：
1. 臺灣糖業史維基百科
2. vocus 糖都之父
3. 楊彥騏《雲林糖業的興衰》
4. 台糖官網虎尾糖廠

四個結構化問題依序問——這是 NotebookLM 的真正用法（Lv.2），不是只問一句就完事。

每問完一題，NotebookLM 會給有出處的答案。提醒學生：不一個一個複製，最後一次性打包（下一張）。""")


# Step 1 - 打包指令
slide = content_slide("1", "最後一步：請它幫你打包", ACCENT)
add_text(slide, Inches(0.8), Inches(1.85), Inches(11.5), Inches(0.5),
         "問完所有問題後，貼這段指令：", size=20, color=LIGHT_GRAY)
add_callout(slide, Inches(0.8), Inches(2.5), Inches(11.5), Inches(2.7),
            "請把我們剛才所有對話整合成一份「研究筆記」，格式：\n\n"
            "一、研究問題     二、關鍵發現（每點標明出處）\n"
            "三、比較分析     四、建議大綱\n"
            "五、參考資料清單",
            color=ACCENT, size=15)
add_callout(slide, Inches(0.8), Inches(5.6), Inches(11.5), Inches(0.85),
            "▸  這份研究筆記就是接下來寫小論文的全部素材",
            color=ACCENT3, size=17)

add_notes(slide, """關鍵：問完所有問題後，不要一個一個複製。

貼這段「整合成研究筆記」的指令，NotebookLM 會把前面所有零碎答案整合成一份結構化筆記，每一點都標明出處。

下一張就是匯出這份筆記成 Word 檔，餵給 Claude 寫小論文。""")


# Step 1 - 匯出成 Word
slide = content_slide("1", "匯出成 Word 檔", ACCENT)
add_text(slide, Inches(0.8), Inches(1.85), Inches(11.5), Inches(0.6),
         "NotebookLM 不適合生成正式文件，但「文獻整合 + 出處標註」很強",
         size=18, color=LIGHT_GRAY)

# 兩步驟卡片
steps12 = [
    ("1. 點「儲存到記事」", "在剛才的回答底下", ACCENT),
    ("2. 工作室 → 匯出至文件", "拿到一份 .docx Word 檔", ACCENT3),
]
for i, (title, desc, color) in enumerate(steps12):
    x = Inches(0.8 + i * 6.0)
    y = Inches(2.8)
    add_rect(slide, x, y, Inches(5.7), Inches(2.8), CARD_BG)
    add_rect(slide, x, y, Inches(5.7), Pt(5), color)
    add_text(slide, x + Inches(0.3), y + Inches(0.4), Inches(5.3), Inches(0.6),
             title, size=22, bold=True, color=WHITE, font=FONT_TITLE)
    add_text(slide, x + Inches(0.3), y + Inches(1.4), Inches(5.3), Inches(1.0),
             desc, size=16, color=LIGHT_GRAY)

add_callout(slide, Inches(0.8), Inches(6.0), Inches(11.5), Inches(1.0),
            "拿到 .docx → 下一步上傳給 Claude 寫小論文",
            color=ACCENT4, size=18)

add_notes(slide, """現場示範這兩步操作（教學網頁有截圖）：

1. 在 NotebookLM 對話的回答底下點「儲存到記事」——記事會出現在「工作室」面板。
2. 點開記事 → 選「匯出至文件」→ 下載一份 .docx Word 檔。

這份 Word 檔包含完整的研究筆記與出處，下一步 Claude 直接吃這個檔當素材，比複製貼上乾淨多了。""")


# ═══════════════════════════════════════════
# Step 2: Napkin
# ═══════════════════════════════════════════
slide = step_divider(2, "Napkin", "把文字變成概念結構圖", ACCENT2)
add_notes(slide, """第二段「做」25 min。Napkin 段落只 3 min。

關鍵訊息：好小論文的核心不是文字，是「概念結構」——觀點之間的關係。""")


slide = content_slide("2", "Napkin——畫一張概念結構圖", ACCENT2)
add_text(slide, Inches(0.8), Inches(1.85), Inches(11.5), Inches(0.5),
         "好的小論文核心不只是文字——是「概念結構」",
         size=22, bold=True, color=WHITE, font=FONT_TITLE)
add_text(slide, Inches(0.8), Inches(2.5), Inches(11.5), Inches(0.5),
         "把研究筆記裡的五個關鍵因素貼進 napkin.ai：",
         size=18, color=LIGHT_GRAY)
add_callout(slide, Inches(0.8), Inches(3.2), Inches(11.5), Inches(2.7),
            "日本選擇在雲林虎尾設立糖廠的原因：\n"
            "1. 天災：1897 戊戌大水災沖出大量溪埔地\n"
            "2. 自然：雲林平原日照充足、雨量適中\n"
            "3. 政策：1902 糖業獎勵規則\n"
            "4. 企業家：鈴木藤三郎，1906 年設廠\n"
            "5. 國家需求：日俄戰爭後日本需糖自給自足",
            color=ACCENT2, size=15)
add_callout(slide, Inches(0.8), Inches(6.2), Inches(11.5), Inches(0.85),
            "→ Napkin 自動生成因果圖 → 下載備用，下一步餵給 Claude",
            color=ACCENT3, size=17)

add_notes(slide, """重點：論文不只是「寫出來」，是「想出來」——用圖把背後的結構表達出來。

NotebookLM 整理的五個原因，貼進 napkin.ai，自動生成因果關係圖。下載備用，下一步上傳給 Claude（連同 NotebookLM 的 Word）。

可選：如果學生想得更深，也可以從研究筆記其他角度生圖（例如比較圖、時間軸）。重點是訓練「看出結構」這件事。""")


# ═══════════════════════════════════════════
# Step 3: Claude
# ═══════════════════════════════════════════
slide = step_divider(3, "Claude", "生成小論文 Word 檔", ACCENT3)
add_notes(slide, """Claude 段落 10 min。產出小論文 Word 初稿。""")


# Step 3 - 怎麼用
slide = content_slide("3", "Claude——生成 Word", ACCENT3)
add_text(slide, Inches(0.8), Inches(1.85), Inches(11.5), Inches(0.5),
         "打開 claude.ai（免費帳號），依序操作：",
         size=20, color=LIGHT_GRAY)
order = [
    ("1", "貼上小論文格式規範（見教學網頁附錄）", ACCENT3),
    ("2", "上傳 Step 1 NotebookLM 匯出的 Word 檔", ACCENT3),
    ("3", "上傳 Step 2 Napkin 因果關係圖", ACCENT3),
    ("4", "指令：「整理成小論文，輸出 Word（含篇名、作者、大綱、引註）」", ACCENT3),
]
for i, (num, text, color) in enumerate(order):
    y = Inches(2.6 + i * 0.85)
    badge = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.8), y, Inches(0.5), Inches(0.5))
    badge.fill.solid(); badge.fill.fore_color.rgb = color; badge.line.fill.background()
    tf = badge.text_frame
    tf.margin_left = tf.margin_right = Pt(0); tf.margin_top = tf.margin_bottom = Pt(0)
    tf.paragraphs[0].text = num
    tf.paragraphs[0].font.size = Pt(20); tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    add_text(slide, Inches(1.55), y + Inches(0.1), Inches(11), Inches(0.5),
             text, size=18, color=WHITE)
add_callout(slide, Inches(0.8), Inches(6.2), Inches(11.5), Inches(0.85),
            "▸  Claude 生成 Artifact → 點「下載」→ 拿到 .docx",
            color=ACCENT4, size=17)

add_notes(slide, """現場示範整個流程：

1. 貼格式規範（教學網頁附錄整段複製）。可以一起上傳格式說明 PDF。
2. 上傳 NotebookLM 匯出的 Word（裡面有研究筆記與出處）。
3. 上傳 Napkin 做好的因果關係圖（Claude 看得到圖）。
4. 下指令——指令裡可以給篇名、作者、大綱要點。指令詳見教學網頁。

Claude 生成 Artifact（Word 預覽），點下載拿到 .docx。

備案：Claude 額度用完，用 Gemini，到 Google Docs 後下載成 Word。""")


# Step 3 - 下載後檢查
slide = content_slide("3", "下載後一定要做的事", ACCENT3)
checks = [
    ("換成你的口吻", "教授看得出 AI 味", ACCENT),
    ("檢查事實", "人名、年份、事件——AI 會編", ACCENT4),
    ("反思先留著", "下一步會改", ACCENT3),
]
for i, (title, desc, color) in enumerate(checks):
    x = Inches(0.8 + i * 4.0)
    add_rect(slide, x, Inches(2.5), Inches(3.7), Inches(3.5), CARD_BG)
    add_rect(slide, x, Inches(2.5), Inches(3.7), Pt(4), color)
    add_text(slide, x, Inches(3.0), Inches(3.7), Inches(0.6),
             title, size=24, bold=True, color=WHITE,
             font=FONT_TITLE, align=PP_ALIGN.CENTER)
    add_text(slide, x, Inches(4.0), Inches(3.7), Inches(0.4),
             desc, size=16, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

add_notes(slide, """關鍵觀念：AI 幫你產出初稿，但「你的判斷」才是學習歷程的靈魂。

三件事：
1. 換成自己的口吻——AI 用詞太成熟、太教科書，一看就是 AI 寫的。
2. 檢查事實——AI 真的會編人名、年份、事件，照單全收會出包。
3. 反思先留著——AI 寫的反思是罐頭，下一段（ChatGPT）會用 prompt 三招改。

這份 Word 還不能直接交。""")


# ═══════════════════════════════════════════
# Step 4: ChatGPT
# ═══════════════════════════════════════════
slide = step_divider(4, "ChatGPT", "用 AI 當教練改反思", ACCENT4)
add_notes(slide, """第三段「學」20 min。

學習歷程最難的不是論文本體，是「反思」。流水帳的反思一看就是 AI 寫的。""")


# Step 4 - Prompt 三招
slide = content_slide("4", "Prompt 三招：給角色、給範例、給限制", ACCENT4)
add_callout(slide, Inches(0.8), Inches(1.85), Inches(11.5), Inches(0.85),
            "「幫我寫一段學習歷程反思」  →  罐頭、流水帳，沒有你的影子",
            color=ACCENT2, label="✗  爛 prompt", size=15)
add_callout(slide, Inches(0.8), Inches(2.95), Inches(11.5), Inches(2.0),
            "我住雲林，常經過虎尾糖廠，以前只覺得是老地方。研究後才知道日本人選雲林"
            "不是巧合，是一場水災沖出了空地。幫我用「從不在意到驚訝」的轉變寫反思，"
            "像高中生寫的，300 字以內。",
            color=ACCENT3, label="✓  好 prompt", size=14)
# 三招卡片
tricks = [
    ("1. 角色", "住雲林的高中生", ACCENT),
    ("2. 範例", "從「不在意」到「驚訝」", ACCENT4),
    ("3. 限制", "300 字、高中生口吻", ACCENT3),
]
for i, (title, desc, color) in enumerate(tricks):
    x = Inches(0.8 + i * 4.0)
    y = Inches(5.2)
    add_rect(slide, x, y, Inches(3.7), Inches(1.3), CARD_BG)
    add_rect(slide, x, y, Inches(3.7), Pt(4), color)
    add_text(slide, x, y + Inches(0.2), Inches(3.7), Inches(0.5),
             title, size=20, bold=True, color=color,
             font=FONT_TITLE, align=PP_ALIGN.CENTER)
    add_text(slide, x, y + Inches(0.75), Inches(3.7), Inches(0.4),
             desc, size=14, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
add_text(slide, Inches(0.8), Inches(6.7), Inches(11.5), Inches(0.4),
         "你給 AI 多少「你的樣子」，它就回給你多少「你的樣子」",
         size=14, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

add_notes(slide, """對照爛 prompt vs. 好 prompt 的差距。

好 prompt 給 AI 三件事：
1. 角色：住雲林的高中生
2. 範例：從「不在意」到「驚訝」的轉變
3. 限制：300 字、高中生口吻

關鍵句：「你給 AI 多少『你的樣子』，它就回給你多少『你的樣子』。」

學生實作：用 Prompt 三招寫一段自己的反思初稿。""")


# Step 4 - 教練 prompt
slide = content_slide("4", "用 AI 改，不是用 AI 寫", ACCENT4)
add_text(slide, Inches(0.8), Inches(1.85), Inches(11.5), Inches(0.5),
         "拿到初稿後，繼續用 AI 當教練：", size=20, color=LIGHT_GRAY)
prompts = [
    ("這段反思沒有連結到研究問題，幫我指出哪裡可以改。", ACCENT),
    ("給我三種角度：學到什麼 / 對家鄉看法改變 / 還想探索什麼", ACCENT4),
    ("這句話怎麼改更有力？", ACCENT3),
]
for i, (text, color) in enumerate(prompts):
    y = Inches(2.6 + i * 1.05)
    add_callout(slide, Inches(0.8), y, Inches(11.5), Inches(0.85),
                f"「{text}」", color=color, size=16)
add_callout(slide, Inches(0.8), Inches(6.0), Inches(11.5), Inches(0.95),
            "至少來回三輪。AI 是逼你想清楚的教練，不是幫你寫的槍手。",
            color=ACCENT2, size=18)

add_notes(slide, """三個示範問題逐一示範。

學習歷程最難的不是報告本體，是「反思」。AI 是逼你想清楚的教練——不是幫你寫的槍手。

至少來回三輪：寫初稿 → AI 指出問題 → 你自己改 → 再請 AI 看。

學生實作：跟 AI 來回改自己的反思段落，至少三輪。""")


# Step 4 - 貼回 Claude
slide = content_slide("4", "改完反思，貼回 Claude 統一風格", ACCENT4)
add_text(slide, Inches(0.8), Inches(1.85), Inches(11.5), Inches(0.5),
         "反思改好後，回到 Step 3 的 Claude 對話框：", size=20, color=LIGHT_GRAY)
add_callout(slide, Inches(0.8), Inches(2.6), Inches(11.5), Inches(2.5),
            "這是我改寫過的反思段落：[貼上]\n\n"
            "請把它整合進剛才的小論文，並確保整篇風格一致。\n"
            "重新輸出 Word 檔。",
            color=ACCENT, size=17)
add_callout(slide, Inches(0.8), Inches(5.4), Inches(11.5), Inches(1.5),
            "下載最新版 Word——\n"
            "這時候才真正算是「你的」學習歷程",
            color=ACCENT3, label="↓  最後一步", size=18)

add_notes(slide, """關鍵流程：反思改好之後不是貼回 Word 檔自己編輯，而是貼回 Step 3 的 Claude 對話框，讓 Claude 幫你把整篇文章的風格統一。

這樣整篇文章從引言到反思，語氣、用詞會一致——不會出現「論文段落是 Claude 寫的、反思是 ChatGPT 寫的」的拼貼感。

下載最新版的 Word 檔，這時候才真正是「你的」學習歷程——AI 幫你把思考變清晰，但內容是你的。""")


# ═══════════════════════════════════════════
# Step 5: Gamma
# ═══════════════════════════════════════════
slide = step_divider(5, "Gamma", "做簡報", ACCENT5)
add_notes(slide, """第四段「簡報」15 min。""")


slide = content_slide("5", "Gamma——做簡報", ACCENT5)
add_text(slide, Inches(0.8), Inches(1.85), Inches(11.5), Inches(0.5),
         "把 Word 內容貼進 gamma.app，輸入：",
         size=20, color=LIGHT_GRAY)
add_callout(slide, Inches(0.8), Inches(2.6), Inches(11.5), Inches(2.0),
            "高中歷史課程學習成果報告簡報。\n"
            "標題：（你的題目）\n"
            "一、研究動機  二、研究發現  三、學習反思",
            color=ACCENT5, size=16)
add_text(slide, Inches(0.8), Inches(4.9), Inches(11.5), Inches(0.5),
         "→ 30 秒生成簡報", size=20, bold=True, color=ACCENT)
add_text(slide, Inches(0.8), Inches(5.5), Inches(11.5), Inches(0.5),
         "→ 記得刪掉 AI 亂加的內容、補上自己的話",
         size=20, bold=True, color=ACCENT)
add_text(slide, Inches(0.8), Inches(6.3), Inches(11.5), Inches(0.4),
         "也可以用 claude.ai 直接生成 .pptx",
         size=16, color=SUBTITLE_GRAY)

add_notes(slide, """關鍵觀念：簡報是用來「說」的，不是用來「讀」的——每張投影片只留關鍵字，細節用嘴巴講。

把 Word 全文貼進 Gamma，加一段提示詞，30 秒生成簡報。現場快速改一兩張：刪 AI 亂加的、補上自己的話。

備案：Claude 也可以直接生成 .pptx 下載。

學生實作：用 Gamma 或 Claude 把 Word 變簡報，跟旁邊同學互看「哪裡看得出是你做的、哪裡像 AI 罐頭」。""")


# ═══════════════════════════════════════════
# Bonus: AI Agent
# ═══════════════════════════════════════════
slide = section_divider("BONUS", "進階：AI Agent", "讓你的學習歷程跟隔壁不一樣", ACCENT4)
add_notes(slide, """第五段「創造」20 min。重頭戲。

從 prompt-copy-paste 進化到 AI Agent——AI 自己規劃、自己執行。
你的學習歷程不只是文字檔，可以是互動網頁、時間軸、資訊圖表。""")


# Bonus - 互動時間軸
slide = content_slide(None, "互動時間軸網頁", ACCENT)
add_text(slide, Inches(0.8), Inches(1.85), Inches(11.5), Inches(0.5),
         "用 Claude Code 或 Gemini CLI，一句話下指令：",
         size=20, color=LIGHT_GRAY)
add_callout(slide, Inches(0.8), Inches(2.5), Inches(11.5), Inches(1.8),
            "把這份雲林糖業報告做成一個互動時間軸網站：\n"
            "1906 大日本製糖設廠 → 1924 日搾量東洋第一 →\n"
            "1945 戰後台糖接收 → 2000s 轉型文創園區",
            color=ACCENT, size=15)
add_text(slide, Inches(0.8), Inches(4.6), Inches(11.5), Inches(0.5),
         "→ AI 自己建檔案 → 寫 HTML/CSS/JS → 跑起來",
         size=18, color=WHITE)
add_text(slide, Inches(0.8), Inches(5.2), Inches(11.5), Inches(0.5),
         "→ 追加：「加一個三題的選擇題小測驗」→ AI 自己改 code",
         size=18, color=WHITE)
add_callout(slide, Inches(0.8), Inches(6.1), Inches(11.5), Inches(0.85),
            "互動網頁連結放進學習歷程——審查委員會記得你",
            color=ACCENT3, size=17)

add_notes(slide, """現場示範 Claude Code 或 Gemini CLI。

關鍵：學生不需要會寫程式，只需要會「下指令」。AI Agent 自己建檔、寫 code、跑起來，再追加需求 AI 自己改。

「同樣兩小時，你能交出的東西多了三倍。」

提醒紅線：最後一關永遠是你——AI 交成品，掛名的是你。""")


# Bonus - 封面視覺
slide = content_slide(None, "封面視覺：Ideogram", ACCENT4)
add_text(slide, Inches(0.8), Inches(1.85), Inches(11.5), Inches(0.5),
         "用 ideogram.ai 生成封面圖：",
         size=20, color=LIGHT_GRAY)
add_callout(slide, Inches(0.8), Inches(2.5), Inches(11.5), Inches(1.5),
            "「日治時期台灣糖廠，五分車鐵道，甘蔗田，復古海報風格」",
            color=ACCENT4, size=17)
add_text(slide, Inches(0.8), Inches(4.4), Inches(11.5), Inches(0.5),
         "Ideogram 特別擅長有文字標語的海報",
         size=20, color=WHITE)
add_callout(slide, Inches(0.8), Inches(5.5), Inches(11.5), Inches(1.5),
            "最後一關永遠是你——\n"
            "AI 交成品，掛名的是你",
            color=ACCENT2, label="🚩  紅線提醒", size=18)

add_notes(slide, """快速示範 Ideogram。

兩個用法：
1. 學習歷程封面圖：放專題報告首頁，視覺亮眼。
2. 補充 Step 5 簡報的視覺素材。

紅線：最後一關永遠是你——AI 交成品，掛名的是你。

收尾這段：學生實作（選擇）：把前段的糖業史報告選一段做成互動網頁，或用 Ideogram 做學習歷程封面圖。""")


# ═══════════════════════════════════════════
# AI 怎麼辦到的？
# ═══════════════════════════════════════════
slide = section_divider("揭開黑盒子", "AI 怎麼辦到的？", "輸入 → 函數 → 輸出", ACCENT5)
add_notes(slide, """第六段 15 min。揭開黑盒子。

接下來的順序：函數本質 → 學測落點例子 → 28×28 手寫辨識例子 → 總結表 → 接龍函數 → 輪盤/溫度 → Token → 訓練三階段 → 幻覺 → 觀察重點 → 不該丟給 AI。""")


# AI = 函數（基本概念）
slide = content_slide(None, "AI 的本質是「函數」", ACCENT5)
add_text(slide, Inches(0.8), Inches(1.85), Inches(11.5), Inches(0.5),
         "國中數學學過：輸入一個值 → 經過規則 → 得到輸出",
         size=20, color=LIGHT_GRAY)
add_callout(slide, Inches(0.8), Inches(2.7), Inches(11.5), Inches(1.5),
            "華氏溫度  =  1.8  ×  攝氏溫度  +  32",
            color=ACCENT, size=24)
add_text(slide, Inches(0.8), Inches(4.6), Inches(11.5), Inches(0.5),
         "寫成 y = f(x)——黑盒子裡有兩個參數（1.8、32），輸入攝氏，算出華氏",
         size=18, color=LIGHT_GRAY)
add_callout(slide, Inches(0.8), Inches(5.7), Inches(11.5), Inches(1.0),
            "AI 做的事情一模一樣——只是函數複雜得多",
            color=ACCENT3, size=20)

add_notes(slide, """國中數學學過函數：輸入 → 規則 → 輸出。

舉最簡單的例子：華氏 = 1.8 × 攝氏 + 32。寫成 y = f(x)，函數有兩個參數（1.8、32）。

AI 做的事情一模一樣，只是函數複雜得多——下兩張會看更接近 AI 的例子（學測落點、手寫辨識）。""")


# AI 函數例子1: 學測落點
slide = content_slide(None, "例子 1：學測落點預測", ACCENT5)
add_text(slide, Inches(0.8), Inches(1.85), Inches(11.5), Inches(0.5),
         "輸入：學測成績組合     輸出：可能錄取的大學機率",
         size=20, color=LIGHT_GRAY)
add_callout(slide, Inches(0.8), Inches(2.7), Inches(11.5), Inches(2.4),
            "f(x) = 國文 × 0.3\n"
            "        + 英文 × 0.25\n"
            "        + 數學 × 0.35\n"
            "        + 社會 × 0.05\n"
            "        + 自然 × 0.05",
            color=ACCENT4, size=18)
add_text(slide, Inches(0.8), Inches(5.4), Inches(11.5), Inches(0.5),
         "實際上更複雜（志願、歷年分布、加權…），但原理一樣",
         size=16, color=SUBTITLE_GRAY)
add_callout(slide, Inches(0.8), Inches(6.1), Inches(11.5), Inches(0.95),
            "原理：輸入 → 函數 → 輸出",
            color=ACCENT3, size=20)

add_notes(slide, """從學生熟悉的「學測落點預測」切入：那也是 y = f(x)，y 是錄取機率，x 是成績組合。

簡化版函數：把每科加權後相加。實際上更複雜，但原理一樣：輸入 → 函數 → 輸出。

下一張看更接近 AI 本質的例子——手寫辨識。""")


# AI 函數例子2: 28x28
slide = content_slide(None, "例子 2：手寫數字辨識", ACCENT5)
add_text(slide, Inches(0.8), Inches(1.85), Inches(11.5), Inches(0.5),
         "輸入 28×28 個小數（每個像素的亮度），輸出 0–9 的機率：",
         size=18, color=LIGHT_GRAY)

# 嵌入兩張圖
img_pixels = "images/04-digit5-pixels.png"
img_prob = "images/04-digit5-prob.png"
if os.path.exists(img_pixels):
    slide.shapes.add_picture(img_pixels, Inches(0.8), Inches(2.6),
                              width=Inches(5.7), height=Inches(3.4))
if os.path.exists(img_prob):
    slide.shapes.add_picture(img_prob, Inches(6.8), Inches(2.6),
                              width=Inches(5.7), height=Inches(3.4))

add_callout(slide, Inches(0.8), Inches(6.2), Inches(11.5), Inches(0.85),
            "AI 看到的不是「5」，是一堆數字；輸出的不是「答案」，是一堆機率",
            color=ACCENT, size=17)

add_notes(slide, """關鍵概念：AI 看到的不是「圖片」，是 28×28=784 個數字（每個 0-255 的亮度值）。

輸出也不是「答案」，是十個機率（0 的機率、1 的機率…9 的機率）。最高的機率對應「答案」。

這就是手寫辨識函數的本質——輸入向量、輸出向量。其他 AI 任務都類似（貓狗辨識、文字生成、圖片生成）。""")


# AI 函數總結
slide = content_slide(None, "AI 能做什麼？都是函數", ACCENT5)
tasks = [
    ("手寫辨識", "一張圖片", "0–9 各數字的機率"),
    ("貓狗辨識", "一張照片", "貓 90%、狗 10%"),
    ("文字生成", "一段提示詞", "生成的文字"),
    ("影像生成", "一段描述", "生成的圖片"),
]
# Header
add_rect(slide, Inches(0.8), Inches(1.9), Inches(11.5), Inches(0.55), CARD_BG)
add_text(slide, Inches(1.0), Inches(2.0), Inches(2.5), Inches(0.4),
         "任務", size=18, bold=True, color=ACCENT5, font=FONT_TITLE)
add_text(slide, Inches(3.6), Inches(2.0), Inches(3.0), Inches(0.4),
         "輸入", size=18, bold=True, color=ACCENT5, font=FONT_TITLE)
add_text(slide, Inches(6.8), Inches(2.0), Inches(1.2), Inches(0.4),
         "→ f →", size=18, bold=True, color=ACCENT5,
         font=FONT_TITLE, align=PP_ALIGN.CENTER)
add_text(slide, Inches(8.2), Inches(2.0), Inches(4.0), Inches(0.4),
         "輸出", size=18, bold=True, color=ACCENT5, font=FONT_TITLE)
for i, (task, inp, out) in enumerate(tasks):
    y = Inches(2.55 + i * 0.65)
    add_rect(slide, Inches(0.8), y, Inches(11.5), Inches(0.55),
             RGBColor(0x1F, 0x28, 0x40))
    add_text(slide, Inches(1.0), y + Inches(0.13), Inches(2.5), Inches(0.4),
             task, size=16, color=WHITE)
    add_text(slide, Inches(3.6), y + Inches(0.13), Inches(3.0), Inches(0.4),
             inp, size=16, color=LIGHT_GRAY)
    add_text(slide, Inches(6.8), y + Inches(0.13), Inches(1.2), Inches(0.4),
             "→ f →", size=14, color=ACCENT5, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(8.2), y + Inches(0.13), Inches(4.0), Inches(0.4),
             out, size=16, color=LIGHT_GRAY)
add_callout(slide, Inches(0.8), Inches(5.8), Inches(11.5), Inches(1.4),
            "不管多厲害，骨子裡都是：輸入 → 函數 → 輸出\n"
            "差別只在於函數裡有幾千億個參數（GPT-4 超過一兆）",
            color=ACCENT3, size=18)

add_notes(slide, """總結 AI 函數家族：手寫辨識、貓狗辨識、文字生成、影像生成——都是函數。

差別只在於函數裡有幾千億個參數（GPT-4 超過一兆個），這些參數從海量資料中「學」出來。

這就是 AI 的本質。下一張開始講「生成式 AI」這種特定類型的函數。""")


# 接龍函數
slide = content_slide(None, "生成式 AI 是「會接龍的函數」", ACCENT5)
add_text(slide, Inches(0.8), Inches(1.85), Inches(11.5), Inches(0.5),
         "每次只猜下一個字，再把輸出接回輸入",
         size=22, color=LIGHT_GRAY)
add_callout(slide, Inches(0.8), Inches(2.7), Inches(11.5), Inches(3.0),
            'f("什麼是臺灣最美的風景？")           →  臺\n'
            'f("什麼是臺灣最美的風景？臺")         →  灣\n'
            'f("什麼是臺灣最美的風景？臺灣")       →  最\n'
            'f("什麼是臺灣最美的風景？臺灣最")     →  美\n'
            '   ……\n'
            'f("什麼是臺灣最美的風景？臺灣最美的風景是")  →  人',
            color=ACCENT, size=15)
add_text(slide, Inches(0.8), Inches(6.1), Inches(11.5), Inches(0.5),
         "就像 LINE 鍵盤自動建議下一個字——AI 做一樣的事，只是猜得非常準",
         size=18, color=LIGHT_GRAY)

add_notes(slide, """既然 AI 是函數，那生成式 AI 是什麼樣的函數？

答案：每次只猜「下一個字」的函數，把自己的輸出接回輸入，像接龍。

用「什麼是臺灣最美的風景？」逐字示範：第一次預測「臺」，把「臺」接回去再預測「灣」，依此類推，最後產生「臺灣最美的風景是人」。

就像 LINE 鍵盤自動建議下一個字——AI 做一樣的事，只是猜得非常準。""")


# 輪盤 / 溫度
slide = content_slide(None, "為什麼同一題會得到不同答案？", ACCENT5)
add_text(slide, Inches(0.8), Inches(1.85), Inches(11.5), Inches(0.5),
         "AI 不是每次都選機率最高的字",
         size=24, bold=True, color=WHITE, font=FONT_TITLE)
add_text(slide, Inches(0.8), Inches(2.55), Inches(11.5), Inches(0.5),
         "像轉「輪盤」——機率高的容易中，但不是一定",
         size=20, color=LIGHT_GRAY)
add_callout(slide, Inches(0.8), Inches(3.5), Inches(11.5), Inches(2.5),
            "溫度低 → 回答比較固定\n\n"
            "溫度高 → 回答有創意，但比較容易出錯",
            color=ACCENT4, label="🌡  溫度參數（Temperature）控制隨機性",
            size=18)
add_text(slide, Inches(0.8), Inches(6.4), Inches(11.5), Inches(0.5),
         "→ 這個「猜字 + 輪盤」也是幻覺的根源（下張揭曉）",
         size=16, color=ACCENT, bold=True)

add_notes(slide, """關鍵概念：AI 不是每次選機率最高的字，加入了一點隨機性（溫度參數）。

像轉輪盤：機率高的字佔輪盤多大區域，但每次轉到的不一定。
- 溫度低：選擇集中、回答比較固定（適合事實性問題）
- 溫度高：選擇分散、回答比較有創意（適合創作）

這也是為什麼問同一題兩次可能得到不同答案。

這個機制下一張會跟「幻覺」連結——AI 不是在「回憶事實」，是在「猜下一個最可能的字」。""")


# Token
slide = content_slide(None, "Token：AI 眼中的世界", ACCENT5)
add_text(slide, Inches(0.8), Inches(1.85), Inches(11.5), Inches(0.5),
         "AI 看到的不是文字，是 Token（碎片）",
         size=24, bold=True, color=WHITE, font=FONT_TITLE)
add_callout(slide, Inches(0.8), Inches(2.7), Inches(11.5), Inches(2.5),
            "▸  輸入「虎尾糖廠」  →  看它被切成幾個 token\n"
            "▸  輸入「314159」  →  數字的切法跟你想的完全不一樣",
            color=ACCENT, label="打開 platform.openai.com/tokenizer 試試",
            size=17)
add_text(slide, Inches(0.8), Inches(5.7), Inches(11.5), Inches(0.5),
         "這就是為什麼 AI 有時候數學會算錯——它看到的不是數字，是碎片",
         size=18, color=ACCENT2, bold=True)

add_notes(slide, """為什麼 AI 數學會算錯？因為它看到的不是「數字」，是 Token（碎片）。

現場演示 OpenAI Tokenizer：
- 「虎尾糖廠」中文怎麼切（中文斷詞特別有趣）
- 「314159」數字怎麼切——可能不是 3-1-4-1-5-9，而是 31-41-59 之類的

這就是為什麼 AI 算 314159 × 217 會算錯——它看到的是奇怪的碎片，不是數字。""")


# 訓練三階段
slide = content_slide(None, "AI 怎麼被訓練出來的？", ACCENT5)
stages = [
    ("1", "預訓練", "讀了整個網路的文字\n→ 學會語言模式", ACCENT),
    ("2", "微調", "針對任務練習\n→ 學會當助理", ACCENT4),
    ("3", "RLHF", "人類回饋好壞\n→ 學會避免有害回應", ACCENT3),
]
for i, (num, title, desc, color) in enumerate(stages):
    x = Inches(0.8 + i * 4.0)
    add_rect(slide, x, Inches(2.3), Inches(3.7), Inches(4.2), CARD_BG)
    add_rect(slide, x, Inches(2.3), Inches(3.7), Pt(4), color)
    add_text(slide, x, Inches(2.7), Inches(3.7), Inches(1.2),
             num, size=56, bold=True, color=color,
             font=FONT_TITLE, align=PP_ALIGN.CENTER)
    add_text(slide, x, Inches(4.3), Inches(3.7), Inches(0.6),
             title, size=24, bold=True, color=WHITE,
             font=FONT_TITLE, align=PP_ALIGN.CENTER)
    add_text(slide, x, Inches(5.1), Inches(3.7), Inches(1.2),
             desc, size=15, color=LIGHT_GRAY,
             align=PP_ALIGN.CENTER, line_spacing=1.5)

add_notes(slide, """三階段：
1. 預訓練：餵整個網路的文字，學語言模式（最耗運算的階段）
2. 微調（Fine-tuning）：針對「當助理」這個任務專門練習
3. RLHF（Reinforcement Learning from Human Feedback）：人類告訴它什麼是好答案、什麼是壞答案，學會避免有害回應

這是 ChatGPT、Claude 之類的 chat 模型怎麼煉出來的。""")


# 幻覺
slide = content_slide(None, "幻覺——AI 也會騙你", ACCENT2)
add_text(slide, Inches(0.8), Inches(1.85), Inches(11.5), Inches(0.5),
         "AI 有時候會一本正經地「編故事」，這叫幻覺",
         size=22, bold=True, color=WHITE, font=FONT_TITLE)
add_callout(slide, Inches(0.8), Inches(2.7), Inches(11.5), Inches(2.0),
            "還記得剛才的「輪盤」嗎？AI 不是在「回憶事實」，是在「猜下一個最可能的字」"
            "——它沒辦法分辨「真的」跟「聽起來像真的」",
            color=ACCENT4, label="❓  為什麼會有幻覺？", size=15)
add_callout(slide, Inches(0.8), Inches(5.0), Inches(11.5), Inches(2.0),
            "分別打開 ChatGPT、Claude、Gemini、Grok，貼上完全相同的問題：\n"
            "「請描述 1920 年雲林農民組合抗議虎尾糖廠壓低甘蔗收購價格的經過」",
            color=ACCENT2, label="實驗：用同一個問題測四個 AI", size=15)

add_notes(slide, """關鍵連結：把幻覺接回前面講的「輪盤+猜字」機制。

AI 不是在「回憶事實」（它沒有資料庫），是在「猜下一個最可能的字」。如果它沒看過這個事實，但「聽起來像真的」的字機率高，它就會吐出來——這就是幻覺。

實驗：用一個 AI 不太熟的台灣史問題（1920 雲林農民組合抗議），測四家 AI，看誰會編、誰會老實說「不確定」。""")


# 觀察重點
slide = content_slide(None, "觀察重點", ACCENT2)
add_bullets(slide, Inches(0.8), Inches(1.95), Inches(11.5), Inches(2.5), [
    "四個 AI 講的是同一件事嗎？",
    "有沒有 AI 承認「我不確定」？",
    "跟 NotebookLM 的素材比對——哪些是真的？哪些是編的？",
], size=20, bullet_color=ACCENT2)
add_callout(slide, Inches(0.8), Inches(4.7), Inches(11.5), Inches(2.0),
            "AI 不是知識庫，是語言模型。\n"
            "所以才需要 NotebookLM——有出處的回答才可信。",
            color=ACCENT3, label="結論", size=20)

add_notes(slide, """觀察四家 AI 的差異：通常會看到「四個 AI 給你四個版本的歷史」。

關鍵警告：「大學教授一年看幾千份學習歷程，你掰的他看得出來。」

結論：AI 不是知識庫，是語言模型——所以才需要 NotebookLM 這種「有餵資料、有出處」的工具。

這也是把整場工作坊收回主線：為什麼一定要從 NotebookLM 開始而不是 ChatGPT。""")


# 不該丟給 AI
slide = content_slide(None, "什麼東西不該丟給 AI？", ACCENT2)
add_text(slide, Inches(0.8), Inches(1.85), Inches(11.5), Inches(0.5),
         "你打進去的每一個字，都會離開你的電腦",
         size=22, color=LIGHT_GRAY)
add_callout(slide, Inches(0.8), Inches(2.7), Inches(5.6), Inches(3.5),
            "▸  日記、心情、家裡狀況\n\n"
            "▸  同學個資、輔導紀錄\n\n"
            "▸  未公開的競賽作品",
            color=ACCENT2, label="不要丟", size=17)
add_callout(slide, Inches(6.7), Inches(2.7), Inches(5.6), Inches(3.5),
            "「這段話被截圖貼到網路上，\n我會崩潰嗎？」\n\n"
            "會  →  別丟雲端",
            color=ACCENT3, label="判準", size=17)

add_notes(slide, """另一條紅線：什麼東西不該丟雲端。

理論上，你打進去的每個字都離開了你的電腦——進了 OpenAI/Anthropic/Google 的伺服器。
- 不要丟：日記、感情、家裡狀況、同學個資、輔導紀錄、未公開的競賽作品
- 進階解法（點到為止）：本機 AI（如 ollama + 開源模型）跑在自己筆電上，資料不外流

判準：「這段話如果被截圖貼到網路上，我會崩潰嗎？」會 → 別丟雲端。""")


# ═══════════════════════════════════════════
# 收尾
# ═══════════════════════════════════════════

# 你現在手上有什麼？
slide = content_slide(None, "你現在手上有什麼？", ACCENT)
products = [
    ("讀", "NotebookLM 研究筆記（Word 匯出）", ACCENT),
    ("做", "Napkin 因果圖 + Claude 小論文 Word 初稿", ACCENT2),
    ("學", "來回三輪改出來的反思，更新進 Word 完成定稿", ACCENT3),
    ("簡報", "Gamma 或 Claude 生成的報告簡報", ACCENT4),
    ("創造", "（加分）互動時間軸網頁 + Ideogram 封面圖", ACCENT5),
]
for i, (stage, what, color) in enumerate(products):
    y = Inches(1.95 + i * 0.78)
    add_rect(slide, Inches(0.8), y, Inches(11.5), Inches(0.65), CARD_BG)
    add_rect(slide, Inches(0.8), y, Pt(5), Inches(0.65), color)
    add_text(slide, Inches(1.0), y + Inches(0.13), Inches(1.0), Inches(0.4),
             stage, size=20, bold=True, color=color, font=FONT_TITLE)
    add_text(slide, Inches(2.2), y + Inches(0.15), Inches(10.0), Inches(0.4),
             what, size=18, color=WHITE)
add_callout(slide, Inches(0.8), Inches(6.1), Inches(11.5), Inches(0.95),
            "三小時前你還在煩惱學習歷程——現在你手上有 Word + 簡報，可能還有互動網頁",
            color=ACCENT3, size=17)

add_notes(slide, """把今天做的東西全部攤開——全部圍繞同一主題：「為什麼是雲林？」

三小時前學生還在煩惱學習歷程怎麼寫，現在手上有 Word + 簡報，可能還有互動網頁。

這就是 Lv.3 的成果——三小時用 AI 做出別人三週才做得出的成品。""")


# 老師示範作品
slide = content_slide(None, "老師示範：跑完流程長這樣", ACCENT3)
add_text(slide, Inches(0.8), Inches(1.85), Inches(11.5), Inches(0.5),
         "以「為什麼是雲林？」為題跑完整套流程，產出：",
         size=20, color=LIGHT_GRAY)
demos = [
    ("⏱  互動時間軸", "虎尾糖廠時間軸.html", ACCENT),
    ("📄  小論文 Word 檔", "Step 3 + 4 定稿", ACCENT4),
    ("🎞  Gamma 簡報", "Step 5 路徑 A", ACCENT3),
    ("🎞  Claude 簡報", "Step 5 路徑 B", ACCENT5),
]
for i, (title, desc, color) in enumerate(demos):
    col = i % 2
    row = i // 2
    x = Inches(0.8 + col * 6.0)
    y = Inches(2.6 + row * 1.7)
    add_rect(slide, x, y, Inches(5.7), Inches(1.5), CARD_BG)
    add_rect(slide, x, y, Inches(5.7), Pt(4), color)
    add_text(slide, x + Inches(0.3), y + Inches(0.25), Inches(5.3), Inches(0.5),
             title, size=22, bold=True, color=WHITE, font=FONT_TITLE)
    add_text(slide, x + Inches(0.3), y + Inches(0.85), Inches(5.3), Inches(0.5),
             desc, size=15, color=LIGHT_GRAY)
add_callout(slide, Inches(0.8), Inches(6.3), Inches(11.5), Inches(0.7),
            "三小時後，你手上也會有類似的東西——而且是你自己選的題目、你自己的口吻",
            color=ACCENT3, size=15)

add_notes(slide, """揭曉老師事先用「為什麼是雲林？」這題跑完整套流程的成品：
- 互動時間軸（HTML）：用 Claude Code 做的，可以放進學習歷程「自主學習」
- 小論文 Word：Step 3+4 定稿
- Gamma 簡報：Step 5 路徑 A，視覺強
- Claude 簡報：Step 5 路徑 B，可對照兩種風格

教學策略選擇：可以開場就揭露當動機；或保留到結尾再揭露當對照。

學生可從工作坊網頁下載這四份示範檔案參考。""")


# 三個帶走觀念
slide = content_slide(None, "三個帶走的觀念", ACCENT)
takeaways = [
    ("1", "AI 幫你加速，不是幫你代寫",
     "最後掛名的是你", ACCENT),
    ("2", "好 prompt 比好工具重要",
     "你給 AI 多少「你的樣子」，它就回給你多少", ACCENT4),
    ("3", "一定要自己檢查",
     "四個 AI 可能給你四個版本的「歷史」", ACCENT3),
]
for i, (num, title, desc, color) in enumerate(takeaways):
    y = Inches(2.0 + i * 1.55)
    add_rect(slide, Inches(0.8), y, Inches(11.5), Inches(1.3), CARD_BG)
    add_rect(slide, Inches(0.8), y, Pt(5), Inches(1.3), color)
    badge = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(1.1), y + Inches(0.3), Inches(0.7), Inches(0.7))
    badge.fill.solid(); badge.fill.fore_color.rgb = color; badge.line.fill.background()
    tf = badge.text_frame
    tf.margin_left = tf.margin_right = Pt(0); tf.margin_top = tf.margin_bottom = Pt(0)
    tf.paragraphs[0].text = num
    tf.paragraphs[0].font.size = Pt(28); tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    add_text(slide, Inches(2.1), y + Inches(0.2), Inches(10), Inches(0.6),
             title, size=24, bold=True, color=WHITE, font=FONT_TITLE)
    add_text(slide, Inches(2.1), y + Inches(0.75), Inches(10), Inches(0.4),
             desc, size=16, color=LIGHT_GRAY)

add_notes(slide, """收尾的三個觀念：

1. AI 幫你加速，不是幫你代寫——掛名的是你（責任感）
2. 好 prompt 比好工具重要——你給 AI 多少「你的樣子」（駕馭能力）
3. 一定要自己檢查——四個 AI 可能給你四個版本的歷史（批判思考）

下一步：今天教的是 Lv.1→Lv.3，Lv.4 正在發生（2026 起）。
回家作業（給有興趣的人）：去看一下 Anthropic Skills、ChatGPT Custom GPT，挑一個試試看。""")


# 謝謝
slide = new_slide(BG_DARK)
add_rect(slide, Inches(5.67), Inches(2.8), Inches(2), Pt(3), ACCENT)
add_text(slide, Inches(0), Inches(3.0), SLIDE_W, Inches(1.4),
         "謝謝！", size=72, bold=True, color=WHITE,
         font=FONT_TITLE, align=PP_ALIGN.CENTER)
add_text(slide, Inches(0), Inches(4.5), SLIDE_W, Inches(0.6),
         "現在開始跟 AI 作朋友吧",
         size=24, color=ACCENT, align=PP_ALIGN.CENTER)
add_text(slide, Inches(0), Inches(5.8), SLIDE_W, Inches(0.5),
         "教學網頁：letranger.github.io/2026-04-30-AI-work/",
         size=16, color=SUBTITLE_GRAY, align=PP_ALIGN.CENTER)
add_notes(slide, """收尾。鼓勵學生：今天學的不是「終點」，是「起點」。

最後留 5 分鐘給 Q&A，並提醒：教學網頁有完整文字、素材清單、附錄、四份示範作品下載。""")


# ═══════════════════════════════════════════
output_path = os.path.join(os.path.dirname(__file__), "和AI作朋友-工作坊簡報.pptx")
prs.save(output_path)
print(f"✅ 簡報已儲存：{output_path}")
print(f"📊 共 {len(prs.slides)} 張投影片（含演講者備忘稿）")
